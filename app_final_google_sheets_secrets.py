import streamlit as st

# 🚀 Must be the very first Streamlit command!
st.set_page_config(
    page_title="Data Analyzer",
    layout="wide",
    initial_sidebar_state="expanded"
)

import gspread
from oauth2client.service_account import ServiceAccountCredentials
import bcrypt
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import io
import datetime
import requests
import json

st.markdown("<style>footer{visibility:hidden;}</style>", unsafe_allow_html=True)
st.title("📊 Data Analyzer")

try:
    scope = [
        "https://spreadsheets.google.com/feeds",
        "https://www.googleapis.com/auth/drive"
    ]
    creds_dict = dict(st.secrets["google_sheets"])
    creds = ServiceAccountCredentials.from_json_keyfile_dict(creds_dict, scope)
    client = gspread.authorize(creds)

    auth_sheet = client.open("user_database").worksheet("users")
    history_sheet = client.open("user_database").worksheet("upload_history")
except Exception as e:
    st.error("Google Sheets credentials not found or invalid. Please check .streamlit/secrets.toml.")
    st.stop()

ADMIN_USERNAME = "admin"

def get_users():
    return auth_sheet.get_all_records()

def find_user(username):
    users = get_users()
    for user in users:
        if user["username"] == username:
            return user
    return None

def add_user(username, password):
    hashed_pw = bcrypt.hashpw(password.encode(), bcrypt.gensalt()).decode()
    auth_sheet.append_row([username, hashed_pw])

def delete_user(username_to_delete):
    try:
        data = auth_sheet.get_all_values()
        headers = data[0]
        remaining = [row for row in data if row[0] != username_to_delete and row[0] != "username"]
        auth_sheet.clear()
        auth_sheet.append_row(headers)
        for row in remaining[1:]:
            auth_sheet.append_row(row)
        return True
    except Exception as e:
        st.error(f"Error deleting user: {e}")
        return False

def authenticate(username, password):
    user = find_user(username)
    if user and bcrypt.checkpw(password.encode(), user["password"].encode()):
        return True
    return False

def save_upload_history(username, filename):
    timestamp = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    history_sheet.append_row([username, filename, timestamp])

def get_upload_history():
    return history_sheet.get_all_records()

import openai

def summarize_csv(df):
    openai.api_key = st.secrets["together"][""]  

    openai.api_base = "https://api.together.xyz/v1"

    # Use first few rows as sample content
    preview = df.head(10).to_csv(index=False)
    prompt = (
        "You are a helpful assistant. Summarize this dataset as if for a project report. "
        "Mention number of rows/columns, types of data, and what it looks like.\n\n"
        f"{preview}"
    )

    try:
        response = openai.ChatCompletion.create(
            model="togethercomputer/llama-2-70b-chat",
            messages=[
                {"role": "system", "content": "You are a data analysis assistant."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=300,
            temperature=0.7
        )
        summary = response.choices[0].message.content
        return summary
    except Exception as e:
        return f"Summary generation failed: {e}"


def export_to_ppt(charts, summary, selected_charts):
    prs = Presentation()
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = "Data Analysis Report"
    slide.placeholders[1].text = "Generated via Streamlit"

    if summary:
        bullet_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = "Summary"
        tf = slide.placeholders[1].text_frame
        for sentence in summary.split("."):
            sentence = sentence.strip()
            if sentence:
                tf.add_paragraph().text = sentence

    for title in selected_charts:
        if title in charts:
            fig = charts[title]
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            img_stream = io.BytesIO()
            fig.savefig(img_stream, format="png")
            img_stream.seek(0)
            slide.shapes.add_picture(img_stream, Inches(1), Inches(1.5), width=Inches(8))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

def generate_selected_charts(df, selected_charts, params):
    charts = {}
    numeric_cols = df.select_dtypes(include="number").columns.tolist()

    if "Scatter Plot" in selected_charts:
        fig, ax = plt.subplots()
        sns.scatterplot(data=df, x=params["Scatter Plot"]["x"], y=params["Scatter Plot"]["y"], ax=ax)
        ax.set_title("Scatter Plot")
        charts["Scatter Plot"] = fig
        st.pyplot(fig)

    if "Line Plot" in selected_charts:
        fig, ax = plt.subplots()
        df.plot(x=params["Line Plot"]["x"], y=params["Line Plot"]["y"], ax=ax)
        ax.set_title("Line Plot")
        charts["Line Plot"] = fig
        st.pyplot(fig)

    if "Histogram" in selected_charts:
        fig, ax = plt.subplots()
        df[numeric_cols].hist(ax=ax)
        plt.tight_layout()
        charts["Histogram"] = fig
        st.pyplot(fig)

    if "Box Plot" in selected_charts:
        fig, ax = plt.subplots()
        sns.boxplot(data=df[numeric_cols], ax=ax)
        ax.set_title("Box Plot")
        charts["Box Plot"] = fig
        st.pyplot(fig)

    if "Violin Plot" in selected_charts:
        fig, ax = plt.subplots()
        sns.violinplot(data=df[numeric_cols], ax=ax)
        ax.set_title("Violin Plot")
        charts["Violin Plot"] = fig
        st.pyplot(fig)

    if "Heatmap" in selected_charts:
        fig, ax = plt.subplots()
        sns.heatmap(df[numeric_cols].corr(), annot=True, cmap="coolwarm", ax=ax)
        ax.set_title("Correlation Heatmap")
        charts["Heatmap"] = fig
        st.pyplot(fig)

    if "Pie Chart" in selected_charts:
        col = params["Pie Chart"]["col"]
        counts = df[col].value_counts()
        fig, ax = plt.subplots()
        ax.pie(counts, labels=counts.index, autopct="%1.1f%%", startangle=90)
        ax.axis("equal")
        ax.set_title(f"Pie Chart of {col}")
        charts["Pie Chart"] = fig
        st.pyplot(fig)

    return charts

def main():
    if "logged_in" not in st.session_state:
        st.session_state.logged_in = False
        st.session_state.username = ""

    if not st.session_state.logged_in:
        st.subheader("🔐 Welcome")
        auth_action = st.radio("Choose action", ["Login", "Sign Up"], horizontal=True)
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")

        if auth_action == "Login":
            if st.button("Login"):
                if authenticate(username, password):
                    st.session_state.logged_in = True
                    st.session_state.username = username
                    st.rerun()
                else:
                    st.error("Invalid credentials.")
        else:
            if st.button("Sign Up"):
                if find_user(username):
                    st.error("Username already exists.")
                elif not username or not password:
                    st.error("Provide both username and password.")
                else:
                    add_user(username, password)
                    st.success("✅ Account created! Please Log In.")
        return

    st.sidebar.header("⚙️ Admin Panel")
    st.sidebar.markdown(
        f"Logged in as: <span style='color:lime'>{st.session_state.username}</span>",
        unsafe_allow_html=True
    )
    if st.sidebar.button("Logout"):
        st.session_state.logged_in = False
        st.rerun()

    uploaded_file = st.file_uploader("📁 Upload CSV", type="csv")
    if uploaded_file:
        df = pd.read_csv(uploaded_file)
        st.dataframe(df)
        save_upload_history(st.session_state.username, uploaded_file.name)

        st.subheader("🔍 Filter Data")
        col = st.selectbox("Filter column", df.columns)
        val = st.text_input("Search keyword")
        filtered = df[df[col].astype(str).str.contains(val, na=False)] if val else df
        st.dataframe(filtered)

        st.subheader("📊 Chart Builder")
        numeric_cols = df.select_dtypes(include="number").columns.tolist()
        cat_cols = df.select_dtypes(include="object").columns.tolist()

        available_charts = ["Scatter Plot", "Line Plot", "Histogram", "Box Plot", "Violin Plot", "Heatmap"]
        if cat_cols:
            available_charts.append("Pie Chart")

        selected = st.multiselect("Select charts to include", available_charts)
        chart_params = {}

        if "Scatter Plot" in selected:
            x = st.selectbox("Scatter X", numeric_cols, key="scatter_x")
            y = st.selectbox("Scatter Y", numeric_cols, key="scatter_y")
            chart_params["Scatter Plot"] = {"x": x, "y": y}

        if "Line Plot" in selected:
            x = st.selectbox("Line X", df.columns, key="line_x")
            y = st.multiselect("Line Y", numeric_cols, default=numeric_cols, key="line_y")
            chart_params["Line Plot"] = {"x": x, "y": y}

        if "Pie Chart" in selected:
            cat = st.selectbox("Pie column", cat_cols, key="pie_col")
            chart_params["Pie Chart"] = {"col": cat}
        # Together.ai ChatGPT-style AI Assistant
import openai

openai.api_key = st.secrets["together_api_key"]
openai.api_base = "https://api.together.xyz/v1"

st.subheader("🤖 Ask AI Assistant")
ai_prompt = st.text_area("Ask anything related to data analysis, Python, or your dataset")

if st.button("Ask AI"):
    if not ai_prompt.strip():
        st.warning("Please enter a prompt.")
    else:
        with st.spinner("Thinking..."):
            try:
                response = openai.ChatCompletion.create(
                    model="togethercomputer/llama-2-70b-chat",
                    messages=[
                        {"role": "system", "content": "You are a helpful data analysis assistant."},
                        {"role": "user", "content": ai_prompt}
                    ],
                    max_tokens=300,
                    temperature=0.7
                )
                reply = response["choices"][0]["message"]["content"]
                st.success("💡 AI Response:")
                st.markdown(reply)
            except Exception as e:
                st.error(f"AI failed: {e}")

        if st.button("Export to PPT"):
            charts = generate_selected_charts(df, selected, chart_params)
            summary = summarize_csv(df, token="hf_manideep")
            ppt = export_to_ppt(charts, summary, selected)
            st.download_button("📥 Download PPT", data=ppt, file_name="data_analysis_report.pptx")

    if st.session_state.username == ADMIN_USERNAME:
        st.subheader("🧑‍💼 Admin: Manage Users / History")
        users = get_users()
        deletable = [u["username"] for u in users if u["username"] != ADMIN_USERNAME]
        to_del = st.selectbox("Delete User", deletable)
        if st.button("Delete"):
            if delete_user(to_del):
                st.success(f"User '{to_del}' deleted.")

        st.subheader("📜 Upload History")
        st.dataframe(pd.DataFrame(get_upload_history()))

if __name__ == "__main__":
    main()
